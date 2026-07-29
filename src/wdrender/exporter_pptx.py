# 정적 슬라이드 exporter — 씬 스틸 캡처(image=풀블리드 / hybrid=배경+텍스트박스)와 문서 구조(structure: 표지·목차·섹션·본문·요약 네이티브 조립)로 PPTX 를 만든다
from __future__ import annotations

import io
from contextlib import ExitStack
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .config import RenderConfig, apply_format, load_config
from .page_session import RenderSession
from .pptx_text import EMU_PER_POINT, extract_text_boxes, map_font_family, parse_css_color
from .server import StaticServer

_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "start": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "end": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

# ── 문서 구조 (structure) ────────────────────────────────────────────────
# 씬 스틸 1:1 이 아닌 "문서형 덱"을 조립하기 위한 슬라이드 종류.
# scene 을 지정하면 그 씬의 스틸(image/hybrid)을 쓰고, 없으면 네이티브 텍스트로 만든다.
SLIDE_KINDS: tuple[str, ...] = ("cover", "toc", "section", "body", "summary")

# 표지·섹션 구분은 관례상 쪽번호·꼬리말을 달지 않는다.
_NO_CHROME_KINDS = frozenset({"cover"})

# 네이티브 슬라이드 팔레트 — hwax-blue raw.palette 에서 옮긴 값(문서형 기본).
# 캡처 씬과 섞여도 어긋나지 않게 같은 색을 쓴다. native_style 로 항목별 덮어쓸 수 있다.
NATIVE_STYLE: dict[str, str] = {
    "page": "#FFFFFF",     # 본문·목차·요약 배경
    "band": "#EBEEFA",     # 섹션 구분 배경 (blueSoft)
    "accent": "#1428A0",   # 표지 배경·강조 (blue)
    "ink": "#101B3E",      # 본문 글자
    "sub": "#57607A",      # 보조 글자
    "line": "#E2E6F0",     # 괘선
    "onAccent": "#FFFFFF",  # 표지 글자
    "font": "맑은 고딕",     # PPTX 지정 폰트 (임베딩 불가 — pptx-hybrid.md §4)
}

# 타이포 배율은 **짧은 변** 기준이다 — 가로(1920×1080)·4:3(1440×1080)·A4(1240×1754)를
# 같은 자로 재기 위해서다(pptx_text.classify_role 과 동일한 정규화 규약).
_TYPE_SCALE: dict[str, float] = {
    "cover_title": 0.085,
    "cover_sub": 0.030,
    "toc_head": 0.050,
    "toc_item": 0.030,
    "section_index": 0.100,
    "section_title": 0.062,
    "body_title": 0.046,
    "body_item": 0.030,
    "chrome": 0.020,       # 쪽번호·꼬리말
}


def _hex_rgb(value: str) -> RGBColor:
    """"#RRGGBB" → RGBColor (앞의 # 는 있어도 없어도 된다)."""
    s = value.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _set_font_name(font, name: str) -> None:
    """latin 과 함께 동아시아(ea) 타입페이스도 지정 — 한글이 대체 폰트로 흘러가지 않게."""
    font.name = name
    rPr = font._element
    latin = rPr.find(qn("a:latin"))
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        latin.addnext(ea)
    ea.set("typeface", name)


def _add_text_box(slide, box: dict[str, Any], epx: float, epy: float, idx: int) -> None:
    """스테이지 px 좌표의 텍스트 상자 하나를 슬라이드 위 네이티브 텍스트박스로 얹는다."""
    tb = slide.shapes.add_textbox(
        Emu(int(round(box["x"] * epx))),
        Emu(int(round(box["y"] * epy))),
        Emu(max(1, int(round(box["w"] * epx)))),
        Emu(max(1, int(round(box["h"] * epy)))),
    )
    tb.name = f"wda-{box['role']}-{idx}"
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 한 줄짜리 설계 텍스트는 폰트 대체로 폭이 늘어나도 접히지 않게 줄바꿈을 끈다.
    # 원래 여러 줄인 텍스트만 설계 폭 안에서 접는다.
    tf.word_wrap = box["lines"] > 1

    p = tf.paragraphs[0]
    p.alignment = _ALIGN.get(box["text_align"], PP_ALIGN.LEFT)
    if box.get("line_height_px"):
        p.line_spacing = Pt(box["line_height_px"] * epy / EMU_PER_POINT)
    font_name = map_font_family(box["font_family"])
    size_pt = box["font_size_px"] * epx / EMU_PER_POINT
    for run in box["runs"]:
        if run.get("br"):
            p.add_line_break()  # <br>·pre 개행 → a:br (python-pptx 의 text 로는 '\v')
            continue
        r = p.add_run()
        r.text = run["text"]
        r.font.size = Pt(size_pt)
        r.font.bold = int(run.get("weight") or 400) >= 600
        r.font.italic = bool(run.get("italic"))
        _set_font_name(r.font, font_name)
        rgb = parse_css_color(run.get("color") or box["color"])
        if rgb:
            r.font.color.rgb = RGBColor(*rgb)


class _Canvas:
    """스테이지 px 로 배치하고 EMU 로 찍는 네이티브 슬라이드 캔버스.

    좌표계는 캡처 슬라이드와 완전히 같다 — 같은 epx/epy 비례 상수를 쓰므로
    캡처 씬과 네이티브 슬라이드를 한 덱에 섞어도 여백선이 어긋나지 않는다.
    """

    def __init__(self, prs, cfg: RenderConfig, style: dict[str, str]) -> None:
        self.prs = prs
        self.w = cfg.width
        self.h = cfg.height
        self.epx = prs.slide_width / cfg.width
        self.epy = prs.slide_height / cfg.height
        self.style = style
        self.unit = min(self.w, self.h)          # 타이포 배율 기준 = 짧은 변
        self.mx = int(round(self.w * 0.062))     # 좌우 여백
        self.my = int(round(self.h * 0.075))     # 상하 여백
        self.font = style["font"]

    def size(self, key: str) -> float:
        """타이포 역할 → 스테이지 px 폰트 크기."""
        return _TYPE_SCALE[key] * self.unit

    @property
    def content_w(self) -> int:
        return self.w - 2 * self.mx

    def background(self, slide, color: str) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_rgb(color)

    def rect(self, slide, x: float, y: float, w: float, h: float, color: str,
             name: str = "wda-rule") -> None:
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(round(x * self.epx))), Emu(int(round(y * self.epy))),
            Emu(max(1, int(round(w * self.epx)))), Emu(max(1, int(round(h * self.epy)))),
        )
        shp.name = name
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex_rgb(color)
        shp.line.fill.background()
        shp.shadow.inherit = False

    def text(self, slide, text: str, *, x: float, y: float, w: float, h: float,
             size_px: float, color: str, bold: bool = False, align: str = "left",
             anchor=MSO_ANCHOR.TOP, wrap: bool = True, name: str = "wda-native") -> None:
        tb = slide.shapes.add_textbox(
            Emu(int(round(x * self.epx))), Emu(int(round(y * self.epy))),
            Emu(max(1, int(round(w * self.epx)))), Emu(max(1, int(round(h * self.epy)))),
        )
        tb.name = name
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = anchor
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = _ALIGN.get(align, PP_ALIGN.LEFT)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size_px * self.epx / EMU_PER_POINT)
        r.font.bold = bold
        _set_font_name(r.font, self.font)
        r.font.color.rgb = _hex_rgb(color)


def _as_entry(raw: Any) -> dict:
    """목차 항목 정규화 — "문자열" 또는 {"text":…, "slide":…} 를 같은 모양으로."""
    if isinstance(raw, dict):
        return {"text": str(raw.get("text", "")), "slide": raw.get("slide")}
    return {"text": str(raw), "slide": None}


def normalize_structure(structure: list[dict]) -> list[dict]:
    """문서 구조를 검증하고 슬라이드 번호·목차 항목을 확정한다.

    한 항목 = 슬라이드 1장이다(같은 씬을 두 번 쓰고 싶으면 항목을 두 개 적는다).
    목차 항목(`items`)을 안 준 toc 는 **section 제목**에서 유도하고, section 이
    하나도 없으면 제목 있는 body 에서 유도한다. 번호는 유도된 자리의 슬라이드 번호다.
    """
    if not structure:
        raise ValueError("structure 가 비었다 — 슬라이드 항목을 1개 이상 넣어라")
    plan: list[dict] = []
    for i, raw in enumerate(structure, 1):
        if not isinstance(raw, dict):
            raise ValueError(f"structure[{i - 1}] 가 매핑이 아니다 — {{'kind': ...}} 형태로 적어라")
        kind = str(raw.get("kind") or "body")
        if kind not in SLIDE_KINDS:
            raise ValueError(
                f"structure[{i - 1}] 의 kind {kind!r} 를 모른다 — 가능한 값은 {list(SLIDE_KINDS)} 다"
            )
        plan.append({
            "kind": kind,
            "slide": i,
            "scene": raw.get("scene"),
            "t": raw.get("t"),
            "title": str(raw.get("title") or ""),
            "subtitle": str(raw.get("subtitle") or ""),
            "items": [_as_entry(x) for x in (raw.get("items") or [])],
            "notes": raw.get("notes"),
        })

    # 섹션 일련번호 (섹션 구분 슬라이드의 큰 숫자)
    n_sec = 0
    for it in plan:
        if it["kind"] == "section":
            n_sec += 1
            it["index"] = n_sec

    derived = [
        {"text": it["title"] or str(it["scene"] or ""), "slide": it["slide"]}
        for it in plan if it["kind"] == "section"
    ] or [
        {"text": it["title"], "slide": it["slide"]}
        for it in plan if it["kind"] == "body" and it["title"]
    ]
    for it in plan:
        if it["kind"] == "toc" and not it["items"]:
            it["items"] = [dict(e) for e in derived]
    return plan


def _native_slide(slide, item: dict, canvas: _Canvas) -> None:
    """씬 캡처 없이 텍스트만으로 슬라이드 한 장을 조립한다."""
    st = canvas.style
    kind = item["kind"]
    w, h, mx, my = canvas.w, canvas.h, canvas.mx, canvas.my
    cw = canvas.content_w

    if kind == "cover":
        canvas.background(slide, st["accent"])
        rule_y = int(round(h * 0.34))
        canvas.rect(slide, mx, rule_y, int(round(cw * 0.14)), max(4, int(round(h * 0.007))),
                    st["onAccent"], name="wda-cover-rule")
        ts = canvas.size("cover_title")
        canvas.text(slide, item["title"], x=mx, y=rule_y + ts * 0.75, w=cw, h=ts * 2.6,
                    size_px=ts, color=st["onAccent"], bold=True, name="wda-cover-title")
        if item["subtitle"]:
            ss = canvas.size("cover_sub")
            canvas.text(slide, item["subtitle"], x=mx, y=rule_y + ts * 3.6, w=cw, h=ss * 2.4,
                        size_px=ss, color=st["onAccent"], name="wda-cover-sub")
        return

    if kind == "section":
        canvas.background(slide, st["band"])
        ns = canvas.size("section_index")
        ts = canvas.size("section_title")
        block_h = ns * 1.2 + ts * 1.6
        top = (h - block_h) / 2
        canvas.text(slide, f"{item.get('index', 1):02d}", x=mx, y=top, w=cw, h=ns * 1.2,
                    size_px=ns, color=st["accent"], bold=True, name="wda-section-index")
        canvas.text(slide, item["title"], x=mx, y=top + ns * 1.25, w=cw, h=ts * 1.6,
                    size_px=ts, color=st["ink"], bold=True, name="wda-section-title")
        return

    # toc / body / summary — 제목 + 괘선 + 항목 목록
    canvas.background(slide, st["page"])
    hs = canvas.size("toc_head") if kind == "toc" else canvas.size("body_title")
    head = item["title"] or ("목차" if kind == "toc" else "")
    canvas.text(slide, head, x=mx, y=my, w=cw, h=hs * 1.4,
                size_px=hs, color=st["ink"], bold=True, name=f"wda-{kind}-title")
    rule_y = my + hs * 1.65
    canvas.rect(slide, mx, rule_y, cw, max(2, int(round(h * 0.003))),
                st["accent"] if kind == "summary" else st["line"], name="wda-head-rule")
    if item["subtitle"]:
        ss = canvas.size("chrome") * 1.4
        canvas.text(slide, item["subtitle"], x=mx, y=rule_y + ss * 0.6, w=cw, h=ss * 1.6,
                    size_px=ss, color=st["sub"], name=f"wda-{kind}-sub")
        rule_y += ss * 2.4

    entries = item["items"]
    if not entries:
        return
    fs = canvas.size("toc_item") if kind == "toc" else canvas.size("body_item")
    top = rule_y + fs * 1.6
    avail = (h - my) - top
    # 항목이 많으면 줄 간격을 줄여 항상 한 화면 안에 들어오게 한다(문서형은 밀도가 생명).
    pitch = min(fs * 2.0, avail / max(1, len(entries)))
    numw = fs * 3.0
    for i, e in enumerate(entries):
        y = top + i * pitch
        if kind == "toc":
            canvas.text(slide, f"{i + 1:02d}", x=mx, y=y, w=numw, h=pitch,
                        size_px=fs, color=st["accent"], bold=True, name=f"wda-toc-no-{i}")
            canvas.text(slide, e["text"], x=mx + numw, y=y, w=cw - numw * 2, h=pitch,
                        size_px=fs, color=st["ink"], name=f"wda-toc-item-{i}")
            if e.get("slide") is not None:
                canvas.text(slide, str(e["slide"]), x=w - mx - numw, y=y, w=numw, h=pitch,
                            size_px=fs, color=st["sub"], align="right",
                            name=f"wda-toc-page-{i}")
        else:
            canvas.rect(slide, mx, y + fs * 0.42, fs * 0.30, fs * 0.30,
                        st["accent"], name=f"wda-bullet-{i}")
            canvas.text(slide, e["text"], x=mx + fs * 0.85, y=y, w=cw - fs * 0.85, h=pitch,
                        size_px=fs, color=st["ink"], name=f"wda-{kind}-item-{i}")


def _add_chrome(slide, canvas: _Canvas, *, number: int | None, total: int, footer: str) -> None:
    """쪽번호(우하)·꼬리말(좌하)을 얹는다 — 문서 제목·날짜는 footer 문자열로 받는다."""
    st = canvas.style
    fs = canvas.size("chrome")
    y = canvas.h - canvas.my * 0.62
    if footer:
        canvas.text(slide, footer, x=canvas.mx, y=y, w=canvas.content_w * 0.7, h=fs * 1.6,
                    size_px=fs, color=st["sub"], wrap=False, name="wda-footer")
    if number is not None:
        canvas.text(slide, f"{number} / {total}", x=canvas.w - canvas.mx - fs * 7,
                    y=y, w=fs * 7, h=fs * 1.6, size_px=fs, color=st["sub"],
                    align="right", wrap=False, name="wda-pageno")


def export_pptx(
    root_dir: str | Path,
    page_relpath: str | Path,
    out_path: str | Path,
    *,
    config: RenderConfig | None = None,
    format_id: str | None = None,
    resources: dict[str, str] | None = None,
    stills: dict[str, list[float]] | None = None,
    notes: dict[str, str] | None = None,
    mode: str = "image",
    structure: list[dict] | None = None,
    slide_numbers: bool = False,
    footer: str = "",
    native_style: dict[str, str] | None = None,
    log: Callable[[str], None] = print,
) -> dict:
    """씬별 정지 화면을 캡처해 PPTX로 조립한다.

    format_id — 포맷 id. 주면 무대 크기와 슬라이드 크기(세로 포맷이면 세로 슬라이드)를
                포맷 스펙에서 가져온다. 미지정 시 render.toml 기본(16:9).
    stills — {씬 name: [씬 로컬 시각(초), ...]} 오버라이드. 미지정 씬은
             progress = default_still_progress(기본 0.9) 시점 1장.
    notes  — {씬 name: 발표자 노트 텍스트}. 씬의 모든 슬라이드에 동일 삽입.
    mode   — "image"(기본, 종전 동작): 화면 그대로 풀블리드 그림 1장.
             "hybrid": 텍스트를 숨긴 배경 캡처 위에 원 좌표의 네이티브 텍스트박스를 얹어
             파워포인트에서 글자를 편집할 수 있게 한다.

    structure — **문서형 덱 조립**. 미지정(None)이면 종전과 완전히 같은 동작
                (씬 × stills 1:1). 지정하면 항목 하나가 슬라이드 한 장이 된다.
                [{kind: "cover"|"toc"|"section"|"body"|"summary",
                  scene?, t?, title?, subtitle?, items?, notes?}]
                `scene` 이 있으면 그 씬 스틸(mode 에 따라 image/hybrid)을 쓰고,
                없으면 네이티브 텍스트 슬라이드를 만든다 — 목차·섹션 구분은 이미지로
                구울 이유가 없다. 씬 참조가 하나도 없으면 브라우저를 아예 띄우지 않는다.
    slide_numbers / footer — 쪽번호("N / 전체")·꼬리말(문서 제목·날짜) 옵션.
                structure 모드에서만 동작하고 표지에는 붙지 않는다.
    native_style — NATIVE_STYLE 항목별 덮어쓰기(색·폰트).

    반환: {slides, scenes, out, mode} 요약 dict (hybrid 는 text_boxes·skipped,
          structure 모드는 native_slides·captured_slides 추가).
    """
    if mode not in ("image", "hybrid"):
        raise ValueError(f"mode 는 'image' 또는 'hybrid' — 받은 값 {mode!r}")
    cfg = apply_format(config or load_config(), format_id)
    stills = stills or {}
    notes = notes or {}
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    plan = normalize_structure(structure) if structure is not None else None

    prs = Presentation()
    prs.slide_width = Emu(cfg.slide_w_emu)
    prs.slide_height = Emu(cfg.slide_h_emu)
    blank_layout = prs.slide_layouts[6]

    # 씬 캡처가 필요한가 — 네이티브 전용 구조면 서버·브라우저를 띄우지 않는다.
    needs_page = plan is None or any(it["scene"] for it in plan)
    canvas = _Canvas(prs, cfg, {**NATIVE_STYLE, **(native_style or {})})
    scenes: list[dict] = []
    n_slides = 0
    n_boxes = 0
    n_native = 0
    skipped: list[dict] = []

    with ExitStack() as stack:
        sess = None
        if needs_page:
            srv = stack.enter_context(StaticServer(root_dir))
            url = srv.url_for(page_relpath)
            log(f"[export_pptx] 페이지 로드 {url}")
            sess = stack.enter_context(RenderSession(
                url,
                width=cfg.width,
                height=cfg.height,
                viewport_margin=cfg.viewport_margin,
                resources=resources,
            ))
            scenes = sess.scenes()
            if not scenes:
                raise RuntimeError("window.OM_SCENES 를 읽을 수 없음 — dc 엔트리가 아님?")
        # 스테이지 px → 슬라이드 EMU 비례 상수 (RenderSession 이 원척 고정을 보장)
        epx, epy = canvas.epx, canvas.epy

        def _capture_slide(name: str, at: float, note: str | None) -> tuple[Any, int]:
            """씬 스틸 1장을 슬라이드로 굽는다 (image/hybrid 공통 경로)."""
            boxes: list[dict] = []
            if mode == "hybrid":
                # seek 은 extract 안에서 — 뽑기와 숨기기를 같은 통과에서 확정한다
                boxes = extract_text_boxes(sess, at, hide=True)
                skipped.extend(
                    dict(s, scene=name, t=round(at, 3)) for s in sess.last_text_skips
                )
            else:
                sess.seek(at)
            png = sess.capture()
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                io.BytesIO(png), 0, 0,
                width=prs.slide_width, height=prs.slide_height,
            )
            for i, box in enumerate(boxes):
                _add_text_box(slide, box, epx, epy, i)
            if note is not None:
                slide.notes_slide.notes_text_frame.text = note
            return slide, len(boxes)

        if plan is None:
            start = 0.0
            for sc in scenes:
                name = sc["name"]
                dur = float(sc["dur"])
                local_times = stills.get(name) or [dur * cfg.default_still_progress]
                for lt in local_times:
                    if not (0.0 <= lt <= dur):
                        raise ValueError(
                            f"씬 '{name}' still {lt}s 가 [0, {dur}] 범위 밖"
                        )
                    _, nb = _capture_slide(name, start + lt, notes.get(name))
                    n_boxes += nb
                    n_slides += 1
                    tail = f" 텍스트 {nb}개" if mode == "hybrid" else ""
                    log(
                        f"[export_pptx] 씬 '{name}' t={start + lt:.2f}s "
                        f"→ 슬라이드 {n_slides}{tail}"
                    )
                start += dur
        else:
            # 씬 이름 → (전역 시작 시각, 길이). 구조 항목의 scene 참조를 시각으로 바꾼다.
            spans: dict[str, tuple[float, float]] = {}
            acc = 0.0
            for sc in scenes:
                spans[sc["name"]] = (acc, float(sc["dur"]))
                acc += float(sc["dur"])
            total = len(plan)
            for it in plan:
                name = it["scene"]
                if name is None:
                    slide = prs.slides.add_slide(blank_layout)
                    _native_slide(slide, it, canvas)
                    if it["notes"] is not None:
                        slide.notes_slide.notes_text_frame.text = str(it["notes"])
                    n_native += 1
                    log(f"[export_pptx] {it['kind']} 네이티브 → 슬라이드 {it['slide']}/{total}")
                else:
                    if name not in spans:
                        raise ValueError(
                            f"structure[{it['slide'] - 1}] 의 씬 '{name}' 이 페이지에 없다 — "
                            f"알려진 씬은 {list(spans)} 다"
                        )
                    start, dur = spans[name]
                    lt = it["t"]
                    if lt is None:
                        lt = (stills.get(name) or [dur * cfg.default_still_progress])[0]
                    lt = float(lt)
                    if not (0.0 <= lt <= dur):
                        raise ValueError(f"씬 '{name}' still {lt}s 가 [0, {dur}] 범위 밖")
                    note = it["notes"] if it["notes"] is not None else notes.get(name)
                    slide, nb = _capture_slide(name, start + lt, note)
                    n_boxes += nb
                    tail = f" 텍스트 {nb}개" if mode == "hybrid" else ""
                    log(
                        f"[export_pptx] {it['kind']} 씬 '{name}' t={start + lt:.2f}s "
                        f"→ 슬라이드 {it['slide']}/{total}{tail}"
                    )
                if it["kind"] not in _NO_CHROME_KINDS and (slide_numbers or footer):
                    _add_chrome(
                        slide, canvas,
                        number=it["slide"] if slide_numbers else None,
                        total=total, footer=footer,
                    )
                n_slides += 1

    prs.save(out_path)
    log(f"[export_pptx] 완료 {out_path} ({n_slides}장, mode={mode})")
    info = {"slides": n_slides, "scenes": len(scenes), "out": str(out_path), "mode": mode}
    if mode == "hybrid":
        info["text_boxes"] = n_boxes
        info["skipped"] = skipped
    if plan is not None:
        info["native_slides"] = n_native
        info["captured_slides"] = n_slides - n_native
    return info
