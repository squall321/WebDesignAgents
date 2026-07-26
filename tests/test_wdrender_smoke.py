# wdrender 스모크 테스트 — 합성 2씬×2초 엔트리로 영상/PPTX exporter 실구동 검증
from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

import pytest
from pptx import Presentation

from wdrender.config import RenderConfig
from wdrender.exporter_pptx import export_pptx
from wdrender.exporter_video import export_video, verify_seek_determinism
from wdrender.page_session import RenderSession, vendor_resources
from wdrender.server import StaticServer

REPO_ROOT = Path(__file__).resolve().parents[1]
RUNTIME_DIR = REPO_ROOT / "web" / "runtime"
VENDOR_DIR = REPO_ROOT / "web" / "vendor"

# 오프라인 폴백(__resources) 경유로 CDN 없이 도는 합성 엔트리 — 폰트 CDN도 없음
_SMOKE_HTML = """<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<script src="./support.js"></script>
</head>
<body>
<x-dc>
<helmet>
<style>html, body { margin: 0; padding: 0; height: 100%; background: #111; }</style>
<script>window.OM_SCENES = '[{"name":"A","dur":2},{"name":"B","dur":2}]';</script>
<script>window.OM_PLAYBACK = '{"mode":"times","count":1}';</script>
</helmet>
<x-import component-from-global-scope="SmokeVideo" from="./animations-v2.jsx ./scenes.jsx" style="position:fixed;inset:0" hint-size="100%,100%"></x-import>
</x-dc>
</body>
</html>
"""

_SMOKE_SCENES = """/* 스모크 테스트용 2씬 — 씬은 localTime의 순수 함수 (엔진 계약 준수) */
const { SceneStage } = window;

function SceneA({ localTime: t }) {
  return (
    <div style={{ position: 'absolute', inset: 0, background: '#123456', color: '#fff',
                  display: 'flex', alignItems: 'center', justifyContent: 'center',
                  fontSize: 160, fontFamily: 'sans-serif' }}>
      A {Math.floor(t * 10)}
    </div>
  );
}

function SceneB({ localTime: t }) {
  return (
    <div style={{ position: 'absolute', inset: 0, background: '#654321', color: '#fff',
                  display: 'flex', alignItems: 'center', justifyContent: 'center',
                  fontSize: 160, fontFamily: 'sans-serif' }}>
      B {Math.floor(t * 10)}
    </div>
  );
}

function SmokeVideo() {
  return (
    <div style={{ position: 'relative', width: '100vw', height: '100vh', overflow: 'hidden' }}>
      <SceneStage width={1920} height={1080} bg="#000" scenes={window.OM_SCENES}
                  playback={window.OM_PLAYBACK}>
        {{ 'A': SceneA, 'B': SceneB }}
      </SceneStage>
    </div>
  );
}
window.SmokeVideo = SmokeVideo;
"""


@pytest.fixture(scope="module")
def build_dir(tmp_path_factory) -> Path:
    """support.js + 엔진 + 합성 씬 + 로컬 vendor로 구성된 스모크 빌드 디렉터리."""
    d = tmp_path_factory.mktemp("smoke_build")
    shutil.copy2(RUNTIME_DIR / "support.js", d / "support.js")
    shutil.copy2(RUNTIME_DIR / "animations-v2.jsx", d / "animations-v2.jsx")
    (d / "vendor").mkdir()
    for f in ("react.production.min.js", "react-dom.production.min.js", "babel.min.js"):
        shutil.copy2(VENDOR_DIR / f, d / "vendor" / f)
    (d / "smoke.dc.html").write_text(_SMOKE_HTML, encoding="utf-8")
    (d / "scenes.jsx").write_text(_SMOKE_SCENES, encoding="utf-8")
    return d


@pytest.fixture(scope="module")
def cfg() -> RenderConfig:
    # 스모크는 저fps로 빠르게 — 4s × 8fps = 32프레임
    return RenderConfig(fps=8)


def _ffprobe_duration(path: Path) -> float:
    out = subprocess.run(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration",
         "-of", "csv=p=0", str(path)],
        capture_output=True, text=True, check=True,
    ).stdout.strip()
    return float(out)


def test_export_video_smoke(build_dir: Path, cfg: RenderConfig, tmp_path: Path):
    out = tmp_path / "smoke.mp4"
    info = export_video(
        build_dir, "smoke.dc.html", out,
        config=cfg, resources=vendor_resources("/vendor"), log=lambda m: None,
    )
    assert out.exists() and out.stat().st_size > 0
    assert info["duration"] == pytest.approx(4.0, abs=0.01)
    assert info["frames"] == 32
    assert _ffprobe_duration(out) == pytest.approx(4.0, abs=0.1)


def test_export_pptx_smoke(build_dir: Path, cfg: RenderConfig, tmp_path: Path):
    out = tmp_path / "smoke.pptx"
    info = export_pptx(
        build_dir, "smoke.dc.html", out,
        config=cfg, resources=vendor_resources("/vendor"),
        stills={"B": [0.5, 1.5]},          # 다단 페이즈 오버라이드 경로 검증
        notes={"A": "씬 A 노트"},
        log=lambda m: None,
    )
    assert info["slides"] == 3  # A 기본 1장 + B 오버라이드 2장
    # python-pptx 재열기 검증
    prs = Presentation(str(out))
    assert len(prs.slides) == 3
    assert prs.slides[0].notes_slide.notes_text_frame.text == "씬 A 노트"


def test_seek_determinism_smoke(build_dir: Path, cfg: RenderConfig):
    with StaticServer(build_dir) as srv:
        with RenderSession(
            srv.url_for("smoke.dc.html"),
            width=cfg.width, height=cfg.height,
            viewport_margin=cfg.viewport_margin,
            resources=vendor_resources("/vendor"),
        ) as sess:
            assert sess.duration == pytest.approx(4.0)
            assert sess.sync_seek, "엔진이 data-om-sync-seek 를 광고해야 함"
            results = verify_seek_determinism(sess, [0.5, 2.5])
            assert all(r["match"] for r in results), results


def test_render_determinism_two_full_renders(build_dir: Path, cfg: RenderConfig):
    """M0 결정성 — 동일 입력을 두 번 독립 렌더해 프레임 diff 가 QA 임계값 이내인지.
    실측(demo_sample 2160프레임)은 max_diff_ratio 0.000002 였다. 스모크(32프레임)로 회귀 고정."""
    from wdrender.exporter_video import verify_render_determinism

    r = verify_render_determinism(
        build_dir, "smoke.dc.html",
        config=cfg, resources=vendor_resources("/vendor"),
        channel_tol=8, log=lambda m: None,
    )
    assert r["frame_count_match"], r
    # QAConfig.frame_match_max_ratio=0.02 이내 (실측상 사실상 0)
    assert r["max_diff_ratio"] <= 0.02, r
