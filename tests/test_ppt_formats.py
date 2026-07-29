# PPT 전용 포맷 테스트 — outputs 존중(영상 생략)·slides 분량 규격·native 슬라이드 구조 조립 검증
from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from wdpipeline.format import (
    FormatError,
    FormatSpec,
    format_presets_briefing,
    check_format_templates,
    is_native_tpl,
    list_formats,
    load_format,
    render_targets,
)
from wdrender.config import RenderConfig, inches_to_emu
from wdrender.exporter_pptx import NATIVE_STYLE, export_pptx, normalize_structure

REPO_ROOT = Path(__file__).resolve().parents[1]
RUNTIME_DIR = REPO_ROOT / "web" / "runtime"
VENDOR_DIR = REPO_ROOT / "web" / "vendor"

PPT_FORMATS = ("deck-doc-16x9", "deck-4x3", "print-a4")


# ── 포맷 스펙 ────────────────────────────────────────────────────────────


def test_three_ppt_formats_are_registered():
    ids = list_formats()
    for fid in PPT_FORMATS:
        assert fid in ids


@pytest.mark.parametrize(
    "fid,stage,slide_in",
    [
        ("deck-doc-16x9", (1920, 1080), (13.333, 7.5)),
        ("deck-4x3", (1440, 1080), (10.0, 7.5)),
        ("print-a4", (1240, 1754), (8.27, 11.69)),
    ],
)
def test_ppt_format_stage_and_slide_size(fid: str, stage, slide_in):
    spec = load_format(fid)
    assert (spec.stage.w, spec.stage.h) == stage
    assert (spec.pptx.slide_w_in, spec.pptx.slide_h_in) == slide_in
    assert spec.outputs == ["pptx"]          # 영상 없음 — 이 라운드의 핵심
    assert spec.narration.enabled is False
    assert spec.skeleton == ["cover", "toc", "section", "body", "summary"]
    assert check_format_templates(spec) == []


def test_ppt_formats_measure_length_in_slides_not_seconds():
    """읽는 자료의 분량 정본은 slides 다. duration 은 장당 명목 12초로 파생된다."""
    spec = load_format("deck-doc-16x9")
    assert spec.slides is not None
    assert (spec.slides.min, spec.slides.target, spec.slides.max) == (5, 15, 30)
    # 파생 — 정본이 아니라 공유 스키마(validate_scenario 등)를 위한 채움값
    assert spec.duration.target == pytest.approx(15 * 12.0)
    assert spec.duration.min == pytest.approx(5 * 12.0)
    assert spec.duration.max == pytest.approx(30 * 12.0)


def test_video_format_keeps_duration_as_the_budget():
    """회귀 — 영상 포맷은 slides 없이 duration 만으로 로드된다."""
    spec = load_format("wide-16x9")
    assert spec.slides is None
    assert (spec.duration.min, spec.duration.target, spec.duration.max) == (20, 90, 600)
    assert spec.outputs == ["video", "pptx"]


def _spec_kwargs(**over) -> dict:
    base = dict(
        id="x-fmt", name_ko="시험", stage={"w": 100, "h": 100},
        skeleton=["a"], template_pool={"a": ["native.a"]},
        pptx={"slide_w_in": 10, "slide_h_in": 7.5},
    )
    base.update(over)
    return base


def test_pptx_only_without_slides_is_rejected():
    with pytest.raises(ValueError) as e:
        FormatSpec.model_validate(_spec_kwargs(
            outputs=["pptx"], duration={"target": 60, "min": 10, "max": 90}))
    assert "slides" in str(e.value)


def test_no_duration_and_no_slides_is_rejected():
    with pytest.raises(ValueError) as e:
        FormatSpec.model_validate(_spec_kwargs(outputs=["pptx"]))
    assert "duration 또는 slides" in str(e.value)


@pytest.mark.parametrize("outputs", [[], ["mp4"], ["video", "gif"]])
def test_unknown_or_empty_outputs_rejected(outputs):
    with pytest.raises(ValueError):
        FormatSpec.model_validate(_spec_kwargs(
            outputs=outputs, slides={"target": 5, "min": 1, "max": 9}))


def test_slides_range_must_be_ordered():
    with pytest.raises(ValueError) as e:
        FormatSpec.model_validate(_spec_kwargs(
            outputs=["pptx"], slides={"target": 2, "min": 5, "max": 9}))
    assert "slides 범위 오류" in str(e.value)


# ── native.* 지시자 ──────────────────────────────────────────────────────


def test_native_pool_is_not_a_scene_template():
    """native.* 는 레지스트리 모듈이 아니다 — 템플릿 조회 대상에서 빠져야 한다."""
    spec = load_format("deck-doc-16x9")
    assert is_native_tpl("native.toc") and not is_native_tpl("tpl.doc-toc")
    # 씬 템플릿만 남고 native.* 는 빠진다 (레지스트리·빌드 바인딩이 소비하는 목록)
    assert spec.tpl_ids() == [
        "tpl.doc-cover", "tpl.doc-toc", "tpl.doc-section", "tpl.doc-body", "tpl.doc-summary",
    ]
    assert spec.native_roles() == spec.skeleton      # 전 역할이 네이티브 폴백을 갖는다
    assert spec.primary_tpl("toc") == "tpl.doc-toc"  # 조립 기본은 문서형 씬 템플릿
    assert spec.template_pool["toc"][-1] == "native.toc"


def test_native_only_pool_loads_without_any_scene_template(tmp_path: Path):
    """native.* 만으로 이뤄진 풀도 성립한다 — 씬 템플릿이 없어도 PPT 는 나온다."""
    root = tmp_path / "formats"
    (root / "n-only").mkdir(parents=True)
    (root / "n-only" / "format.yaml").write_text(
        "id: n-only\nname_ko: 네이티브 전용\nstage: { w: 1920, h: 1080 }\n"
        "slides: { target: 5, min: 1, max: 9 }\nskeleton: [cover]\n"
        "template_pool: { cover: [native.cover] }\noutputs: [pptx]\n"
        "pptx: { slide_w_in: 13.333, slide_h_in: 7.5 }\n",
        encoding="utf-8",
    )
    spec = load_format("n-only", formats_root=root)
    assert spec.tpl_ids() == [] and spec.native_roles() == ["cover"]
    assert check_format_templates(spec) == []


def test_registered_doc_templates_must_be_wired_into_the_pool():
    """문서형 씬 템플릿이 레지스트리에 등록되고 이 포맷을 선언하면 풀에 들어와 있어야 한다.

    지금은 native.* 폴백만 있다(모듈 미등록). 등록되는 순간 이 시험이 켜지면서
    '템플릿은 만들었는데 포맷이 안 쓴다'는 어긋남을 잡는다.
    """
    from wdpipeline.format import load_module_registry, load_module_yaml, module_formats

    registry = load_module_registry()
    for fid in PPT_FORMATS:
        spec = load_format(fid)
        pool = {t for role in spec.skeleton for t in spec.template_pool[role]}
        for mid, entry in registry.items():
            if entry.get("type") != "scene-template":
                continue
            module = load_module_yaml(mid)
            if module is None or fid not in module_formats(module):
                continue
            assert mid in pool, f"{mid} 가 {fid} 를 선언했는데 template_pool 에 없다"


def test_video_format_tpl_ids_unchanged():
    """회귀 — 씬 템플릿 포맷의 tpl_ids 는 그대로다."""
    assert load_format("wide-16x9").tpl_ids() == [
        "tpl.opening", "tpl.problem", "tpl.concept", "tpl.process", "tpl.timeline",
        "tpl.differentiator", "tpl.compare", "tpl.proof", "tpl.dataviz", "tpl.closing",
    ]


# ── outputs 존중 (렌더 타깃) ─────────────────────────────────────────────


def test_render_targets_pptx_only_skips_video():
    for fid in PPT_FORMATS:
        assert render_targets(fid) == ["pptx"]


def test_render_targets_video_and_pptx():
    assert render_targets("wide-16x9") == ["video", "pptx"]
    assert render_targets("short-9x16") == ["video", "pptx"]


def test_render_targets_skip_video_override():
    assert render_targets("wide-16x9", skip_video=True) == ["pptx"]
    assert render_targets("deck-4x3", skip_video=True) == ["pptx"]


def test_render_targets_video_only_format(tmp_path: Path):
    """영상 전용 포맷도 성립한다 — 3가지(영상/슬라이드/양쪽)가 다 되어야 한다."""
    root = tmp_path / "formats"
    (root / "v-only").mkdir(parents=True)
    (root / "v-only" / "format.yaml").write_text(
        "id: v-only\n"
        "name_ko: 영상 전용\n"
        "stage: { w: 1920, h: 1080 }\n"
        "duration: { target: 60, min: 20, max: 90 }\n"
        "skeleton: [opening]\n"
        "template_pool: { opening: [tpl.opening] }\n"
        "outputs: [video]\n"
        "pptx: { slide_w_in: 13.333, slide_h_in: 7.5 }\n",
        encoding="utf-8",
    )
    # 템플릿 실재 검사는 끄고 본다 — 이 시험 포맷은 렌더 타깃 판정만이 관심사다
    assert load_format("v-only", formats_root=root, check_templates=False).outputs == ["video"]
    assert render_targets("v-only", formats_root=root) == ["video"]
    assert render_targets("v-only", formats_root=root, skip_video=True) == []


def test_render_targets_unknown_format_falls_back():
    """포맷을 못 읽으면 기존 동작(영상+슬라이드)으로 떨어진다 — 렌더가 스펙 오류로 멈추지 않게."""
    assert render_targets("no-such-format") == ["video", "pptx"]
    assert render_targets(None) == ["video", "pptx"]
    assert render_targets("") == ["video", "pptx"]


def test_briefing_reports_slides_for_ppt_formats():
    text = format_presets_briefing("deck-doc-16x9")
    assert "목표 15장 (허용 5~30)" in text
    assert "산출 pptx" in text
    assert "cover(tpl.doc-cover)" in text
    assert "무대 1920×1080" in text


def test_briefing_still_reports_seconds_for_video_formats():
    text = format_presets_briefing("wide-16x9")
    assert "목표 90초" in text and "산출 video+pptx" in text


# ── 문서 구조 정규화 ─────────────────────────────────────────────────────


def test_structure_numbers_slides_and_derives_toc():
    plan = normalize_structure([
        {"kind": "cover", "title": "표지"},
        {"kind": "toc"},
        {"kind": "section", "title": "배경"},
        {"kind": "body", "title": "현황", "items": ["가", "나"]},
        {"kind": "section", "title": "제안"},
        {"kind": "body", "title": "안"},
        {"kind": "summary", "title": "요약"},
    ])
    assert [it["slide"] for it in plan] == [1, 2, 3, 4, 5, 6, 7]
    toc = plan[1]["items"]
    assert [e["text"] for e in toc] == ["배경", "제안"]
    assert [e["slide"] for e in toc] == [3, 5]          # 섹션 구분 슬라이드 번호
    assert [it.get("index") for it in plan if it["kind"] == "section"] == [1, 2]


def test_structure_toc_falls_back_to_body_titles():
    plan = normalize_structure([
        {"kind": "toc"},
        {"kind": "body", "title": "첫 장"},
        {"kind": "body"},                                # 제목 없음 — 목차에 안 들어간다
        {"kind": "body", "title": "셋째 장"},
    ])
    assert [(e["text"], e["slide"]) for e in plan[0]["items"]] == [("첫 장", 2), ("셋째 장", 4)]


def test_structure_explicit_toc_items_win():
    plan = normalize_structure([
        {"kind": "toc", "items": [{"text": "손으로 쓴 항목", "slide": 9}]},
        {"kind": "section", "title": "무시될 섹션"},
    ])
    assert plan[0]["items"] == [{"text": "손으로 쓴 항목", "slide": 9}]


@pytest.mark.parametrize("bad,msg", [
    ([], "structure 가 비었다"),
    ([{"kind": "appendix"}], "kind 'appendix' 를 모른다"),
    (["cover"], "매핑이 아니다"),
])
def test_structure_rejects_bad_input(bad, msg):
    with pytest.raises(ValueError) as e:
        normalize_structure(bad)
    assert msg in str(e.value)


# ── 네이티브 전용 조립 (브라우저 없이) ───────────────────────────────────


DOC_STRUCTURE = [
    {"kind": "cover", "title": "화학 소재 브리핑", "subtitle": "2026 상반기 심의 결과"},
    {"kind": "toc"},
    {"kind": "section", "title": "배경"},
    {"kind": "body", "title": "현황", "items": ["조각 1", "조각 2", "조각 3"]},
    {"kind": "body", "title": "쟁점", "items": ["쟁점 하나", "쟁점 둘"], "notes": "발표 메모"},
    {"kind": "body", "title": "대안", "items": ["대안 A", "대안 B"]},
    {"kind": "summary", "title": "요약", "items": ["결론 한 줄"]},
]


def _texts(prs) -> list[str]:
    out: list[str] = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text:
                out.append(shp.text_frame.text)
    return out


@pytest.mark.parametrize("fid,size_in", [
    ("deck-doc-16x9", (13.333, 7.5)),
    ("deck-4x3", (10.0, 7.5)),
    ("print-a4", (8.27, 11.69)),
])
def test_native_only_deck_needs_no_browser(tmp_path: Path, fid: str, size_in):
    """씬 참조가 없으면 페이지를 아예 열지 않는다 — root_dir 가 없어도 만들어져야 한다."""
    out = tmp_path / f"{fid}.pptx"
    info = export_pptx(
        tmp_path / "does-not-exist", "index.html", out,
        format_id=fid, structure=DOC_STRUCTURE, slide_numbers=True,
        footer="화학 소재 브리핑 · 2026-07-29", log=lambda s: None,
    )
    assert info["slides"] == 7
    assert info["native_slides"] == 7 and info["captured_slides"] == 0
    assert info["scenes"] == 0

    prs = Presentation(out)
    assert len(prs.slides) == 7
    assert prs.slide_width == inches_to_emu(size_in[0])
    assert prs.slide_height == inches_to_emu(size_in[1])
    # 씬 캡처가 없으니 그림도 없다
    assert sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13) == 0


def test_native_deck_contents_and_numbering(tmp_path: Path):
    out = tmp_path / "doc.pptx"
    export_pptx(tmp_path, "index.html", out, format_id="deck-doc-16x9",
                structure=DOC_STRUCTURE, slide_numbers=True,
                footer="화학 소재 브리핑 · 2026-07-29", log=lambda s: None)
    prs = Presentation(out)

    # 표지 — 쪽번호·꼬리말이 붙지 않는다
    cover = [sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame]
    assert "화학 소재 브리핑" in cover and "2026 상반기 심의 결과" in cover
    assert not any("/" in t for t in cover)

    # 목차 — 섹션 제목에서 유도되고 슬라이드 번호가 붙는다
    toc = [sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame]
    assert "목차" in toc and "배경" in toc
    assert "3" in toc                                    # 섹션 '배경' 은 3번 슬라이드
    assert "2 / 7" in toc                                # 쪽번호
    assert "화학 소재 브리핑 · 2026-07-29" in toc          # 꼬리말

    # 섹션 구분 — 일련번호
    section = [sh.text_frame.text for sh in prs.slides[2].shapes if sh.has_text_frame]
    assert "01" in section and "배경" in section

    # 본문 — 제목 + 항목 전부
    body = [sh.text_frame.text for sh in prs.slides[3].shapes if sh.has_text_frame]
    assert "현황" in body
    for item in ("조각 1", "조각 2", "조각 3"):
        assert item in body

    # 발표자 노트
    assert prs.slides[4].notes_slide.notes_text_frame.text == "발표 메모"
    assert "결론 한 줄" in _texts(prs)


@pytest.mark.parametrize("fid", PPT_FORMATS)
def test_native_shapes_stay_inside_the_slide(tmp_path: Path, fid: str):
    """세 무대(16:9·4:3·A4 세로) 모두에서 도형이 슬라이드 밖으로 나가지 않아야 한다."""
    out = tmp_path / f"{fid}-bounds.pptx"
    export_pptx(tmp_path, "index.html", out, format_id=fid, structure=DOC_STRUCTURE,
                slide_numbers=True, footer="꼬리말", log=lambda s: None)
    prs = Presentation(out)
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            assert sh.left >= 0 and sh.top >= 0, f"{fid} 슬라이드 {i} {sh.name} 음수 좌표"
            assert sh.left + sh.width <= prs.slide_width, f"{fid} 슬라이드 {i} {sh.name} 우측 이탈"
            assert sh.top + sh.height <= prs.slide_height, f"{fid} 슬라이드 {i} {sh.name} 하단 이탈"


def test_slide_numbers_and_footer_are_optional(tmp_path: Path):
    out = tmp_path / "plain.pptx"
    export_pptx(tmp_path, "index.html", out, format_id="deck-4x3",
                structure=DOC_STRUCTURE, log=lambda s: None)
    prs = Presentation(out)
    texts = _texts(prs)
    assert not any(t.endswith(" / 7") for t in texts)
    assert not any("2026-07-29" in t for t in texts)


def test_native_style_override(tmp_path: Path):
    out = tmp_path / "styled.pptx"
    export_pptx(tmp_path, "index.html", out, format_id="print-a4",
                structure=[{"kind": "cover", "title": "표지"}],
                native_style={"accent": "#A8402F"}, log=lambda s: None)
    prs = Presentation(out)
    fill = prs.slides[0].background.fill
    assert str(fill.fore_color.rgb) == "A8402F"
    assert NATIVE_STYLE["accent"] == "#1428A0"           # 원본 상수는 안 바뀐다


# ── 씬 캡처 + 구조 혼합 (브라우저 실구동) ────────────────────────────────


_HTML = """<!DOCTYPE html>
<html><head><meta charset="utf-8"><script src="./support.js"></script></head>
<body><x-dc><helmet>
<style>html, body { margin: 0; padding: 0; height: 100%; background: #fff; }</style>
<script>window.OM_SCENES = '[{"name":"A","dur":2},{"name":"B","dur":2}]';</script>
<script>window.OM_PLAYBACK = '{"mode":"times","count":1}';</script>
</helmet>
<x-import component-from-global-scope="DeckVideo" from="./animations-v2.jsx ./scenes.jsx"
 style="position:fixed;inset:0" hint-size="100%,100%"></x-import>
</x-dc></body></html>
"""

_SCENES = """/* 구조 조립 검증용 2씬 — 정적 배치 */
const { SceneStage } = window;
function mk(label) {
  return function Scene() {
    return (
      <div style={{ position: 'absolute', inset: 0, background: '#ffffff',
                    fontFamily: "'Pretendard Variable', sans-serif" }}>
        <div style={{ position: 'absolute', left: 200, top: 400, width: 1200, fontSize: 96,
                      fontWeight: 800, color: '#101B3E', lineHeight: '120px' }}>{label}</div>
      </div>
    );
  };
}
function DeckVideo() {
  return (
    <div style={{ position: 'relative', width: '100vw', height: '100vh', overflow: 'hidden' }}>
      <SceneStage width={1920} height={1080} bg="#fff" scenes={window.OM_SCENES}
                  playback={window.OM_PLAYBACK}>
        {{ 'A': mk('씬 에이'), 'B': mk('씬 비') }}
      </SceneStage>
    </div>
  );
}
window.DeckVideo = DeckVideo;
"""


@pytest.fixture(scope="module")
def scene_build(tmp_path_factory) -> Path:
    d = tmp_path_factory.mktemp("deck_build")
    shutil.copy2(RUNTIME_DIR / "support.js", d / "support.js")
    shutil.copy2(RUNTIME_DIR / "animations-v2.jsx", d / "animations-v2.jsx")
    (d / "vendor").mkdir()
    for f in ("react.production.min.js", "react-dom.production.min.js", "babel.min.js"):
        shutil.copy2(VENDOR_DIR / f, d / "vendor" / f)
    (d / "deck.dc.html").write_text(_HTML, encoding="utf-8")
    (d / "scenes.jsx").write_text(_SCENES, encoding="utf-8")
    return d


def _pictures(slide) -> int:
    return sum(1 for sh in slide.shapes if sh.shape_type == 13)


def test_structure_mixes_captured_scenes_and_native_slides(scene_build: Path, tmp_path: Path):
    out = tmp_path / "mixed.pptx"
    info = export_pptx(
        scene_build, "deck.dc.html", out, config=RenderConfig(),
        structure=[
            {"kind": "cover", "title": "혼합 덱"},
            {"kind": "toc", "items": ["씬 에이", "씬 비"]},
            {"kind": "body", "scene": "A", "title": "에이"},
            {"kind": "body", "scene": "B", "t": 0.5, "notes": "비 노트"},
            {"kind": "summary", "title": "끝", "items": ["한 줄"]},
        ],
        slide_numbers=True, log=lambda s: None,
    )
    assert info["slides"] == 5
    assert info["native_slides"] == 3 and info["captured_slides"] == 2
    assert info["scenes"] == 2

    prs = Presentation(out)
    assert [_pictures(s) for s in prs.slides] == [0, 0, 1, 1, 0]
    assert prs.slides[3].notes_slide.notes_text_frame.text == "비 노트"
    # 쪽번호는 캡처 슬라이드 위에도 얹힌다 (표지 제외)
    assert "3 / 5" in [sh.text_frame.text for sh in prs.slides[2].shapes if sh.has_text_frame]


def test_structure_unknown_scene_is_rejected(scene_build: Path, tmp_path: Path):
    with pytest.raises(ValueError) as e:
        export_pptx(scene_build, "deck.dc.html", tmp_path / "x.pptx",
                    config=RenderConfig(), structure=[{"kind": "body", "scene": "없는씬"}],
                    log=lambda s: None)
    assert "없는씬" in str(e.value) and "알려진 씬은" in str(e.value)


def test_structure_still_time_out_of_range_is_rejected(scene_build: Path, tmp_path: Path):
    with pytest.raises(ValueError) as e:
        export_pptx(scene_build, "deck.dc.html", tmp_path / "x.pptx",
                    config=RenderConfig(), structure=[{"kind": "body", "scene": "A", "t": 9.0}],
                    log=lambda s: None)
    assert "범위 밖" in str(e.value)


def test_no_structure_keeps_the_old_behaviour(scene_build: Path, tmp_path: Path):
    """회귀 절대 조건 — structure 미지정은 종전 그대로(씬 × stills 1:1, 풀블리드 그림)."""
    out = tmp_path / "legacy.pptx"
    info = export_pptx(scene_build, "deck.dc.html", out, config=RenderConfig(),
                       notes={"A": "에이 노트"}, log=lambda s: None)
    assert info["slides"] == 2 and info["scenes"] == 2
    assert "native_slides" not in info and "captured_slides" not in info
    prs = Presentation(out)
    assert [_pictures(s) for s in prs.slides] == [1, 1]
    # 텍스트 프레임 0 — image 모드는 그림 1장뿐이고 쪽번호도 붙지 않는다
    assert _texts(prs) == []
    assert prs.slides[0].notes_slide.notes_text_frame.text == "에이 노트"


def test_multiple_stills_still_expand_without_structure(scene_build: Path, tmp_path: Path):
    """회귀 — stills 여러 장 지정은 종전대로 씬당 여러 슬라이드가 된다."""
    info = export_pptx(scene_build, "deck.dc.html", tmp_path / "many.pptx",
                       config=RenderConfig(), stills={"A": [0.2, 1.0, 1.8]},
                       log=lambda s: None)
    assert info["slides"] == 4          # A 3장 + B 1장


def test_unknown_format_id_still_raises():
    with pytest.raises(FormatError):
        load_format("deck-doc-16x9-typo")
