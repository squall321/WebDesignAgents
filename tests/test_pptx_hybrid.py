# 하이브리드 PPTX 테스트 — 텍스트 추출 휴리스틱 단위 검증 + 합성 엔트리로 배경/텍스트박스 실구동 검증
from __future__ import annotations

import io
import shutil
from pathlib import Path

import pytest
from PIL import Image, ImageChops
from pptx import Presentation

from wdrender.config import RenderConfig
from wdrender.exporter_pptx import export_pptx
from wdrender.page_session import RenderSession, vendor_resources
from wdrender.pptx_text import (
    DEFAULT_FALLBACK_FONT,
    classify_role,
    extract_text_boxes,
    map_font_family,
    parse_css_color,
)
from wdrender.server import StaticServer

REPO_ROOT = Path(__file__).resolve().parents[1]
RUNTIME_DIR = REPO_ROOT / "web" / "runtime"
VENDOR_DIR = REPO_ROOT / "web" / "vendor"

# 1920×1080 무대. 좌표를 코드에 박아 두고 PPTX 도형 좌표와 대조한다.
_HTML = """<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<script src="./support.js"></script>
</head>
<body>
<x-dc>
<helmet>
<style>html, body { margin: 0; padding: 0; height: 100%; background: #fff; }</style>
<script>window.OM_SCENES = '[{"name":"T","dur":2}]';</script>
<script>window.OM_PLAYBACK = '{"mode":"times","count":1}';</script>
</helmet>
<x-import component-from-global-scope="HybridVideo" from="./animations-v2.jsx ./scenes.jsx" style="position:fixed;inset:0" hint-size="100%,100%"></x-import>
</x-dc>
</body>
</html>
"""

_SCENES = """/* 하이브리드 추출 검증용 1씬 — 모든 요소는 localTime 무관 정적 배치 */
const { SceneStage } = window;

function SceneT() {
  return (
    <div style={{ position: 'absolute', inset: 0, background: '#ffffff',
                  fontFamily: "'Pretendard Variable', Pretendard, sans-serif" }}>
      {/* 제목 — 강조 span 이 섞인 다중 run */}
      <div style={{ position: 'absolute', left: 200, top: 100, width: 1000, height: 120,
                    fontSize: 96, fontWeight: 800, color: '#101B3E', textAlign: 'center',
                    lineHeight: '120px' }}>
        앞말 <span style={{ color: '#1428A0' }}>강조</span> 뒷말
      </div>
      {/* 필 — 패딩·배경·그림자를 가진 텍스트 요소 (콘텐츠 상자 배치 + 배경 보존) */}
      <div style={{ position: 'absolute', left: 300, top: 300, padding: '20px 40px',
                    background: '#1428A0', color: '#fff', fontSize: 40, fontWeight: 700,
                    borderRadius: 999, boxShadow: '0 16px 36px rgba(20,40,160,0.3)' }}>
        필배지
      </div>
      {/* 본문 2줄 — 지정 폭에서 접힌다 */}
      <div style={{ position: 'absolute', left: 200, top: 500, width: 600,
                    fontSize: 32, fontWeight: 400, color: '#57607A', lineHeight: '48px' }}>
        본문은 지정한 폭 안에서 두 줄로 접히도록 충분히 길게 쓴 문장이다
      </div>
      {/* 회전 — 제외 대상 */}
      <div style={{ position: 'absolute', left: 1400, top: 700, fontSize: 40,
                    color: '#A8402F', transform: 'rotate(12deg)' }}>
        회전텍스트
      </div>
      {/* 반투명 — 제외 대상 */}
      <div style={{ position: 'absolute', left: 200, top: 800, fontSize: 30,
                    color: '#667085', opacity: 0.4 }}>
        페이드중
      </div>
      {/* 캡션 — 하단 밴드 */}
      <div style={{ position: 'absolute', left: 0, right: 0, top: 1000, fontSize: 24,
                    color: '#667085', textAlign: 'center' }}>
        하단 캡션
      </div>
      {/* <br> 줄바꿈 — 파워포인트 줄바꿈(a:br)으로 옮겨져야 한다 */}
      <div style={{ position: 'absolute', left: 1200, top: 850, width: 600, fontSize: 28,
                    color: '#101B3E', lineHeight: '40px' }}>
        윗줄<br />아랫줄
      </div>
    </div>
  );
}

function HybridVideo() {
  return (
    <div style={{ position: 'relative', width: '100vw', height: '100vh', overflow: 'hidden' }}>
      <SceneStage width={1920} height={1080} bg="#fff" scenes={window.OM_SCENES}
                  playback={window.OM_PLAYBACK}>
        {{ 'T': SceneT }}
      </SceneStage>
    </div>
  );
}
window.HybridVideo = HybridVideo;
"""

EMU_PER_PX = 6350  # 1920px → 12,192,000 EMU / 1080px → 6,858,000 EMU (양축 동일)


@pytest.fixture(scope="module")
def build_dir(tmp_path_factory) -> Path:
    d = tmp_path_factory.mktemp("hybrid_build")
    shutil.copy2(RUNTIME_DIR / "support.js", d / "support.js")
    shutil.copy2(RUNTIME_DIR / "animations-v2.jsx", d / "animations-v2.jsx")
    (d / "vendor").mkdir()
    for f in ("react.production.min.js", "react-dom.production.min.js", "babel.min.js"):
        shutil.copy2(VENDOR_DIR / f, d / "vendor" / f)
    (d / "hybrid.dc.html").write_text(_HTML, encoding="utf-8")
    (d / "scenes.jsx").write_text(_SCENES, encoding="utf-8")
    return d


@pytest.fixture(scope="module")
def cfg() -> RenderConfig:
    return RenderConfig()


@pytest.fixture(scope="module")
def extracted(build_dir: Path, cfg: RenderConfig) -> dict:
    """세션 1회로 보이는 상태·숨긴 상태 캡처와 상자를 모두 뽑아 모듈 전체가 공유한다."""
    with StaticServer(build_dir) as srv:
        with RenderSession(
            srv.url_for("hybrid.dc.html"), width=cfg.width, height=cfg.height,
            viewport_margin=cfg.viewport_margin, resources=vendor_resources("/vendor"),
        ) as sess:
            boxes = extract_text_boxes(sess, 1.0, hide=False)
            skips = list(sess.last_text_skips)
            full = Image.open(io.BytesIO(sess.capture())).convert("RGB")
            extract_text_boxes(sess, 1.0, hide=True)
            bg = Image.open(io.BytesIO(sess.capture())).convert("RGB")
            # 원복 확인 — 다음 호출 첫머리에서 되돌아와야 한다
            boxes_again = extract_text_boxes(sess, 1.0, hide=False)
    return {"boxes": boxes, "skips": skips, "full": full, "bg": bg,
            "boxes_again": boxes_again, "stage": {"w": cfg.width, "h": cfg.height}}


# ── 순수 함수 단위 ────────────────────────────────────────────────────────


def test_map_font_family_pretendard_falls_back():
    assert map_font_family("'Pretendard Variable', Pretendard, sans-serif") == "맑은 고딕"
    assert map_font_family("Consolas, monospace") == "Consolas"
    assert map_font_family("") == DEFAULT_FALLBACK_FONT
    assert map_font_family("Unknown Face") == DEFAULT_FALLBACK_FONT


def test_parse_css_color():
    assert parse_css_color("rgb(20, 40, 160)") == (20, 40, 160)
    assert parse_css_color("rgba(255, 255, 255, 0.5)") == (255, 255, 255)
    assert parse_css_color("transparent") is None


@pytest.mark.parametrize(
    "stage,fs,weight,y,expect",
    [
        ({"w": 1920, "h": 1080}, 112, 800, 433, "title"),      # 가로 hero
        ({"w": 1920, "h": 1080}, 56, 800, 162, "title"),       # 가로 sectionTitle
        ({"w": 1920, "h": 1080}, 36, 500, 600, "subtitle"),    # 가로 subtitle
        ({"w": 1920, "h": 1080}, 26, 700, 108, "label"),       # 가로 킥커
        ({"w": 1920, "h": 1080}, 26, 400, 454, "body"),        # 가로 본문
        ({"w": 1920, "h": 1080}, 33, 800, 381, "label"),       # 카드 제목
        ({"w": 1920, "h": 1080}, 24, 400, 995, "caption"),     # 푸터
        ({"w": 1080, "h": 1920}, 96, 800, 200, "title"),       # 세로 hero
        ({"w": 1080, "h": 1920}, 32, 700, 116, "label"),       # 세로 킥커
        ({"w": 1080, "h": 1920}, 40, 400, 1198, "body"),       # 세로 본문
        ({"w": 1080, "h": 1920}, 32, 400, 1794, "caption"),    # 세로 푸터
    ],
)
def test_classify_role_matches_design_scale(stage, fs, weight, y, expect):
    box = {"font_size_px": fs, "font_weight": weight, "y": y, "h": fs * 1.2}
    assert classify_role(box, stage) == expect


# ── 추출 휴리스틱 실구동 ─────────────────────────────────────────────────


def test_extract_leaf_text_only_no_duplicates(extracted):
    texts = [b["text"] for b in extracted["boxes"]]
    assert texts == ["앞말 강조 뒷말", "필배지",
                     "본문은 지정한 폭 안에서 두 줄로 접히도록 충분히 길게 쓴 문장이다",
                     "하단 캡션", "윗줄\n아랫줄"]
    # 강조 span 은 부모의 run 으로 흡수 — 독립 상자가 생기면 중복이다
    assert "강조" not in texts


def test_extract_runs_keep_accent_color(extracted):
    title = extracted["boxes"][0]
    assert [r["text"] for r in title["runs"]] == ["앞말 ", "강조", " 뒷말"]
    assert parse_css_color(title["runs"][0]["color"]) == (16, 27, 62)
    assert parse_css_color(title["runs"][1]["color"]) == (20, 40, 160)


def test_extract_geometry_is_content_box(extracted):
    title, pill, body, caption, _br = extracted["boxes"]
    assert (title["x"], title["y"], title["w"], title["h"]) == (200, 100, 1000, 120)
    # 필은 padding 20px 40px 를 뺀 콘텐츠 상자로 잡혀야 한다
    assert pill["x"] == pytest.approx(340, abs=0.5)
    assert pill["y"] == pytest.approx(320, abs=0.5)
    assert pill["font_size_px"] == 40
    assert body["w"] == 600 and body["lines"] == 2
    assert title["lines"] == 1
    assert caption["text_align"] == "center" and caption["w"] == 1920
    # 정규화 좌표
    assert title["fx"] == pytest.approx(200 / 1920)
    assert title["fy"] == pytest.approx(100 / 1080)


def test_extract_roles(extracted):
    assert [b["role"] for b in extracted["boxes"]] == [
        "title", "subtitle", "body", "caption", "body"
    ]


def test_extract_br_becomes_line_break_run(extracted):
    br_box = extracted["boxes"][4]
    assert br_box["text"] == "윗줄\n아랫줄"
    assert [r.get("text") for r in br_box["runs"]] == ["윗줄", None, "아랫줄"]
    assert br_box["runs"][1]["br"] is True
    assert br_box["lines"] == 2


def test_extract_skips_rotated_and_faded(extracted):
    reasons = {s["text"]: s["reason"] for s in extracted["skips"]}
    assert reasons.get("회전텍스트") == "rotate/skew"
    assert reasons.get("페이드중") == "opacity"


def test_hide_removes_glyphs_but_keeps_decoration(extracted):
    """배경 캡처에서 글자는 사라지고 필의 배경·그림자는 남는다."""
    full, bg = extracted["full"], extracted["bg"]
    diff = ImageChops.difference(full, bg).convert("L")
    mask = diff.point(lambda p: 255 if p > 8 else 0)
    changed = sum(mask.histogram()[1:])
    assert changed > 0, "숨김 전후가 완전히 같다 — 텍스트가 지워지지 않았다"

    # 변화는 전부 추출된 상자(글리프 여유 4px) 안이어야 한다 — 상자를 지우고 남는 게 있으면 손실
    for b in extracted["boxes"]:
        mask.paste(0, (int(b["x"]) - 4, int(b["y"]) - 4,
                       int(b["x"] + b["w"]) + 5, int(b["y"] + b["h"]) + 5))
    outside = sum(mask.histogram()[1:])
    assert outside == 0, f"상자 밖 {outside}px 가 함께 사라졌다 (bbox={mask.getbbox()})"

    # 필 배경(파랑)은 배경 캡처에도 그대로 — 중심 픽셀 색으로 확인
    assert bg.getpixel((400, 340)) == (20, 40, 160)
    # 회전·반투명 텍스트는 숨기지 않으므로 배경에 남는다
    assert full.crop((1380, 690, 1620, 780)) == bg.crop((1380, 690, 1620, 780))


def test_hide_is_restored_for_next_call(extracted):
    assert [b["text"] for b in extracted["boxes_again"]] == [
        b["text"] for b in extracted["boxes"]
    ]
    assert extracted["boxes_again"][0]["color"] == extracted["boxes"][0]["color"]


# ── exporter ─────────────────────────────────────────────────────────────


def test_export_hybrid_places_editable_text(build_dir, cfg, tmp_path):
    out = tmp_path / "hybrid.pptx"
    info = export_pptx(
        build_dir, "hybrid.dc.html", out, config=cfg,
        resources=vendor_resources("/vendor"), notes={"T": "노트"},
        mode="hybrid", log=lambda m: None,
    )
    assert info["mode"] == "hybrid"
    assert info["slides"] == 1 and info["text_boxes"] == 5
    assert {s["reason"] for s in info["skipped"]} == {"rotate/skew", "opacity"}

    prs = Presentation(str(out))
    assert prs.slide_width == 12_192_000 and prs.slide_height == 6_858_000
    slide = prs.slides[0]
    frames = [sh for sh in slide.shapes if sh.has_text_frame]
    assert [f.text_frame.text for f in frames] == [
        "앞말 강조 뒷말", "필배지",
        "본문은 지정한 폭 안에서 두 줄로 접히도록 충분히 길게 쓴 문장이다", "하단 캡션",
        "윗줄\v아랫줄",   # python-pptx 는 a:br 을 수직 탭으로 읽는다
    ]
    # 배경 그림 1장 + 텍스트박스 5개
    assert sum(1 for sh in slide.shapes if sh.shape_type == 13) == 1
    assert slide.notes_slide.notes_text_frame.text == "노트"

    title = frames[0]
    assert title.left == 200 * EMU_PER_PX and title.top == 100 * EMU_PER_PX
    assert title.width == 1000 * EMU_PER_PX and title.height == 120 * EMU_PER_PX
    runs = [r for p in title.text_frame.paragraphs for r in p.runs]
    assert [r.text for r in runs] == ["앞말 ", "강조", " 뒷말"]
    assert runs[0].font.size.pt == 48.0  # 96px × 6350/12700
    assert runs[0].font.bold is True
    assert str(runs[0].font.color.rgb) == "101B3E"
    assert str(runs[1].font.color.rgb) == "1428A0"
    assert runs[0].font.name == "맑은 고딕"
    # 동아시아 타입페이스도 같이 지정돼야 한글이 대체 폰트로 흐르지 않는다
    from pptx.oxml.ns import qn
    assert runs[0].font._element.find(qn("a:ea")).get("typeface") == "맑은 고딕"

    # 한 줄 텍스트는 줄바꿈 끔 / 두 줄 텍스트만 설계 폭 안에서 접는다
    assert title.text_frame.word_wrap is False
    assert frames[2].text_frame.word_wrap is True
    assert frames[2].width == 600 * EMU_PER_PX


def test_export_image_mode_unchanged(build_dir, cfg, tmp_path):
    """기본 모드 회귀 — 텍스트 프레임 없이 풀블리드 그림 1장 그대로."""
    out = tmp_path / "image.pptx"
    info = export_pptx(
        build_dir, "hybrid.dc.html", out, config=cfg,
        resources=vendor_resources("/vendor"), notes={"T": "노트"}, log=lambda m: None,
    )
    assert info["mode"] == "image"
    assert info["slides"] == 1 and "text_boxes" not in info
    prs = Presentation(str(out))
    slide = prs.slides[0]
    assert [sh for sh in slide.shapes if sh.has_text_frame] == []
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    assert len(pics) == 1
    assert (pics[0].left, pics[0].top) == (0, 0)
    assert (pics[0].width, pics[0].height) == (prs.slide_width, prs.slide_height)
    assert slide.notes_slide.notes_text_frame.text == "노트"


def test_export_rejects_unknown_mode(build_dir, cfg, tmp_path):
    with pytest.raises(ValueError):
        export_pptx(build_dir, "hybrid.dc.html", tmp_path / "x.pptx",
                    config=cfg, mode="native", log=lambda m: None)
